"""
Renderer — Planning JSON → PPTX
================================
读取 Planning JSON 合同，调用 slide_builders 渲染 PPTX。
数据层和渲染层完全解耦：改 JSON 不用碰渲染代码。
"""

import json
from pathlib import Path

from .design_system import BrandPalette
from .slide_builders import DeckBuilder
from .planning_schema import validate_plan


class PlanRenderer:
    """
    从 Planning JSON 渲染 PPTX。

    用法：
        renderer = PlanRenderer("examples/sample_plan.json")
        renderer.render("output/deck.pptx")
    """

    def __init__(self, plan_path: str):
        self.plan_path = Path(plan_path)
        self.plan = json.loads(self.plan_path.read_text())
        self.brand = self._init_brand()
        self.deck = DeckBuilder(self.brand)

    def _init_brand(self) -> BrandPalette:
        """从 plan metadata 初始化品牌色"""
        meta = self.plan.get("metadata", {})
        brand_cfg = meta.get("brand", {})
        return BrandPalette(
            primary=brand_cfg.get("primary", "#00377B"),
            secondary=brand_cfg.get("secondary", "#009FDB"),
            accent=brand_cfg.get("accent", "#FFD100"),
            font_family=brand_cfg.get("font_family", "Arial"),
        )

    @staticmethod
    def _infer_number_format(chart: dict) -> str:
        """从 y_axis_title 和数据自动推断数据标签格式"""
        y_title = (chart.get("y_axis_title") or "").lower()

        # 百分比
        if any(k in y_title for k in ["%", "率", "ratio", "percent", "share", "penetration", "增长"]):
            return '#,##0"%"'
        # 评分/指数
        if any(k in y_title for k in ["评分", "score", "index", "rating", "分"]):
            return '#,##0.0'
        # 欧元
        if "€" in y_title or "eur" in y_title:
            return '€#,##0'
        # 人民币
        if any(k in y_title for k in ["¥", "rmb", "cny", "元"]):
            return '¥#,##0'
        # 美元（默认货币）
        if any(k in y_title for k in ["$", "usd", "cost", "revenue", "price", "费", "价"]):
            return '$#,##0'
        # 数量
        if any(k in y_title for k in ["count", "number", "用户", "户", "人", "万", "个"]):
            return '#,##0'

        # 检查数据大小推断
        all_values = []
        for s in chart.get("series", []):
            all_values.extend(s.get("values", []))
        if all_values:
            max_val = max(abs(v) for v in all_values if isinstance(v, (int, float)))
            if max_val > 1000:
                return '#,##0'
            elif max_val > 100:
                return '#,##0'

        return '#,##0'  # 安全默认：纯数字，不加货币符号

    def validate(self) -> bool:
        """渲染前校验 Planning JSON"""
        errors = validate_plan(str(self.plan_path))
        criticals = [e for e in errors if e.severity == "CRITICAL"]
        if criticals:
            print(f"❌ 校验失败 — {len(criticals)} CRITICAL errors:")
            for e in criticals:
                print(f"  🔴 {e}")
            return False
        warnings = [e for e in errors if e.severity == "WARNING"]
        if warnings:
            print(f"⚠️ {len(warnings)} warnings (非阻断):")
            for e in warnings:
                print(f"  🟡 {e}")
        return True

    def render(self, output_path: str) -> str:
        """渲染全部页面，输出 PPTX"""
        if not self.validate():
            raise ValueError("Planning JSON 校验失败，无法渲染")

        slides = self.plan.get("slides", [])
        for slide_data in slides:
            template = slide_data.get("template")
            self._render_slide(template, slide_data)

        n = self.deck.save(output_path)
        print(f"✅ 渲染完成: {n} 页 → {output_path}")
        return output_path

    def _render_slide(self, template: str, data: dict):
        """根据模板类型分发渲染"""
        dispatch = {
            "cover":             self._render_cover,
            "section_divider":   self._render_section_divider,
            "executive_summary": self._render_executive_summary,
            "data_story":        self._render_data_story,
            "comparison":        self._render_comparison,
            "framework":         self._render_framework,
            "table":             self._render_table,
            "recommendation":    self._render_recommendation,
            "appendix":          self._render_appendix,
            "closing":           self._render_closing,
        }
        fn = dispatch.get(template)
        if fn:
            fn(data)
        else:
            print(f"  ⚠️ 未知模板 '{template}'，跳过")

    # ═══════════════════════════════════════════
    # 各模板渲染器
    # ═══════════════════════════════════════════

    def _render_cover(self, data: dict):
        self.deck.add_cover(
            title=data.get("title", ""),
            subtitle=data.get("subtitle", ""),
            date=data.get("date", ""),
        )

    def _render_section_divider(self, data: dict):
        self.deck.add_section_divider(
            title=data.get("action_title", data.get("title", "")),
            number=data.get("number"),
        )

    def _render_executive_summary(self, data: dict):
        kpis = []
        for kpi in data.get("kpis", []):
            kpis.append((
                kpi.get("value", ""),
                kpi.get("label", ""),
                kpi.get("detail", ""),
            ))
        self.deck.add_kpi_slide(
            title=data.get("action_title", ""),
            kpis=kpis,
            subtitle=data.get("subtitle"),
            source=data.get("source"),
        )

    def _render_data_story(self, data: dict):
        chart = data.get("chart", {})
        series = []
        for s in chart.get("series", []):
            series.append((s.get("name", ""), s.get("values", [])))

        self.deck.add_chart_slide(
            title=data.get("action_title", ""),
            chart_type=chart.get("type", "bar"),
            categories=chart.get("categories", []),
            series=series,
            subtitle=data.get("subtitle"),
            y_axis_title=chart.get("y_axis_title"),
            insight=data.get("insight"),
            source=data.get("source"),
            number_format=chart.get("number_format", self._infer_number_format(chart)),
        )

    def _render_comparison(self, data: dict):
        """左右对比页 — 新增模板"""
        from pptx.util import Inches, Pt
        from pptx.enum.text import PP_ALIGN
        from pptx.enum.shapes import MSO_SHAPE

        slide = self.deck._blank()
        self.deck._action_title(slide, data.get("action_title", ""), data.get("subtitle"))

        # 左侧
        left_data = data.get("left", {})
        right_data = data.get("right", {})

        for col_idx, col_data in enumerate([left_data, right_data]):
            x = Inches(0.8 + col_idx * 6.0)
            color = self.brand.secondary if col_idx == 0 else self.brand.shade_70

            # 标签栏
            label_bar = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE,
                                                x, Inches(1.5), Inches(5.5), Inches(0.5))
            label_bar.fill.solid()
            label_bar.fill.fore_color.rgb = color
            label_bar.line.fill.background()
            lp = label_bar.text_frame.paragraphs[0]
            lp.text = col_data.get("label", "")
            lp.font.size = Pt(14)
            lp.font.bold = True
            lp.font.color.rgb = self.brand.white
            lp.font.name = self.brand.font_family
            lp.alignment = PP_ALIGN.CENTER

            # 要点列表
            points = col_data.get("points", [])
            for j, point in enumerate(points):
                tx = slide.shapes.add_textbox(x + Inches(0.2), Inches(2.2 + j * 0.55),
                                               Inches(5.1), Inches(0.45))
                tx.text_frame.word_wrap = True
                p = tx.text_frame.paragraphs[0]
                p.text = f"• {point}"
                p.font.size = Pt(11)
                p.font.color.rgb = self.brand.charcoal
                p.font.name = self.brand.font_family

        if data.get("source"):
            self.deck._add_source(slide, data["source"])

    def _render_framework(self, data: dict):
        """2x2 矩阵 / Issue Tree — 新增模板"""
        from pptx.util import Inches, Pt
        from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
        from pptx.enum.shapes import MSO_SHAPE

        slide = self.deck._blank()
        self.deck._action_title(slide, data.get("action_title", ""), data.get("subtitle"))

        framework_type = data.get("framework_type", "matrix_2x2")
        quadrants = data.get("quadrants", [])

        if framework_type == "matrix_2x2" and len(quadrants) >= 4:
            colors = [self.brand.tint_05, self.brand.tint_15,
                      self.brand.light_gray, self.brand.tint_05]
            positions = [
                (Inches(1.5), Inches(1.8)),   # 左上
                (Inches(7.0), Inches(1.8)),   # 右上
                (Inches(1.5), Inches(4.2)),   # 左下
                (Inches(7.0), Inches(4.2)),   # 右下
            ]
            for i, quad in enumerate(quadrants[:4]):
                x, y = positions[i]
                box = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, x, y,
                                              Inches(5.0), Inches(2.2))
                box.fill.solid()
                box.fill.fore_color.rgb = colors[i]
                box.line.color.rgb = self.brand.mid_gray
                box.line.width = Pt(0.5)
                box.text_frame.word_wrap = True

                p = box.text_frame.paragraphs[0]
                p.text = quad.get("title", "")
                p.font.size = Pt(12)
                p.font.bold = True
                p.font.color.rgb = self.brand.primary
                p.font.name = self.brand.font_family

                desc = quad.get("description", "")
                if desc:
                    p2 = box.text_frame.add_paragraph()
                    p2.text = desc
                    p2.font.size = Pt(9)
                    p2.font.color.rgb = self.brand.charcoal
                    p2.font.name = self.brand.font_family
                    p2.space_before = Pt(6)

            # 轴标签
            if data.get("x_axis"):
                tx = slide.shapes.add_textbox(Inches(4.5), Inches(6.6), Inches(4), Inches(0.3))
                p = tx.text_frame.paragraphs[0]
                p.text = data["x_axis"]
                p.font.size = Pt(9)
                p.font.color.rgb = self.brand.mid_gray
                p.alignment = PP_ALIGN.CENTER
            if data.get("y_axis"):
                tx = slide.shapes.add_textbox(Inches(0.3), Inches(3.5), Inches(1), Inches(0.3))
                p = tx.text_frame.paragraphs[0]
                p.text = data["y_axis"]
                p.font.size = Pt(9)
                p.font.color.rgb = self.brand.mid_gray

        if data.get("source"):
            self.deck._add_source(slide, data["source"])

    def _render_table(self, data: dict):
        headers = data.get("headers", [])
        rows = data.get("rows", [])
        self.deck.add_table_slide(
            title=data.get("action_title", ""),
            headers=headers,
            rows=rows,
            subtitle=data.get("subtitle"),
            source=data.get("source"),
        )

    def _render_recommendation(self, data: dict):
        recs = []
        for r in data.get("recommendations", []):
            recs.append((
                r.get("number", ""),
                r.get("title", ""),
                r.get("detail", ""),
            ))
        self.deck.add_recommendation_slide(
            title=data.get("action_title", ""),
            recommendations=recs,
            subtitle=data.get("subtitle"),
        )

    def _render_appendix(self, data: dict):
        """附录页 — 使用 table 模板"""
        headers = data.get("headers", ["Item", "Detail"])
        rows = data.get("rows", [])
        self.deck.add_table_slide(
            title=data.get("action_title", "Appendix"),
            headers=headers,
            rows=rows,
            source=data.get("source"),
        )

    def _render_closing(self, data: dict):
        self.deck.add_closing(
            title=data.get("title", "Thank You"),
            subtitle=data.get("subtitle", ""),
        )


# ====================================================================
# CLI
# ====================================================================
if __name__ == "__main__":
    import sys
    if len(sys.argv) < 3:
        print("Usage: python renderer.py <plan.json> <output.pptx>")
        sys.exit(1)
    renderer = PlanRenderer(sys.argv[1])
    renderer.render(sys.argv[2])
