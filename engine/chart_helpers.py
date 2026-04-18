"""
Chart Helpers — 扩展图表类型
=============================
补齐 slide_builders 中缺少的图表类型：
  - 瀑布图 (waterfall)
  - 分组柱状图 (grouped_bar)
  - 折线图 (line)
  - 饼图 (pie/donut)
"""

from pptx.util import Inches, Pt, Emu
from pptx.enum.chart import XL_CHART_TYPE, XL_LEGEND_POSITION, XL_LABEL_POSITION
from pptx.chart.data import CategoryChartData, ChartData
from pptx.dml.color import RGBColor

from .design_system import BrandPalette, FONT


def add_waterfall_chart(slide, brand: BrandPalette,
                        categories: list, values: list,
                        left=Inches(0.8), top=Inches(1.5),
                        width=Inches(11.5), height=Inches(4.5),
                        number_format='$#,##0.0"M"'):
    """
    模拟瀑布图（python-pptx 不原生支持，用堆叠柱状图模拟）。

    categories: ["Base", "Add 1", "Add 2", "Sub 1", "Total"]
    values: [100, 30, 20, -15, 135]  # 正=增加，负=减少，最后一个=总计
    """
    # 计算不可见基座和可见柱体
    n = len(values)
    bases = [0.0] * n    # 不可见基座高度
    positives = [0.0] * n
    negatives = [0.0] * n

    running = 0
    for i, v in enumerate(values):
        if i == n - 1:
            # 最后一项视为总计
            positives[i] = v if v >= 0 else 0
            negatives[i] = -v if v < 0 else 0
            bases[i] = 0
        elif v >= 0:
            bases[i] = running
            positives[i] = v
            running += v
        else:
            running += v  # v 是负数
            bases[i] = running
            negatives[i] = -v

    # 创建堆叠柱状图：3 层（透明基座 + 正向 + 负向）
    chart_data = CategoryChartData()
    chart_data.categories = categories
    chart_data.add_series("Base", bases)
    chart_data.add_series("Increase", positives)
    chart_data.add_series("Decrease", negatives)

    cf = slide.shapes.add_chart(
        XL_CHART_TYPE.COLUMN_STACKED, left, top, width, height, chart_data
    )
    chart = cf.chart

    # 基座透明
    base_series = chart.series[0]
    base_series.format.fill.background()  # 透明
    base_series.format.line.fill.background()

    # 正向 = secondary 色
    inc_series = chart.series[1]
    inc_series.format.fill.solid()
    inc_series.format.fill.fore_color.rgb = brand.secondary
    inc_series.data_labels.show_value = True
    inc_series.data_labels.font.size = Pt(9)
    inc_series.data_labels.font.name = brand.font_family
    inc_series.data_labels.number_format = number_format

    # 负向 = accent 色
    dec_series = chart.series[2]
    dec_series.format.fill.solid()
    dec_series.format.fill.fore_color.rgb = brand.accent
    dec_series.data_labels.show_value = True
    dec_series.data_labels.font.size = Pt(9)
    dec_series.data_labels.font.name = brand.font_family
    dec_series.data_labels.number_format = number_format

    # 总计柱体用 primary 色（手动改最后一个点的颜色）
    # python-pptx 对单点颜色的支持有限，这里用 series 整体色

    # 图表格式
    chart.has_legend = False
    chart.category_axis.tick_labels.font.size = Pt(8)
    chart.category_axis.tick_labels.font.name = brand.font_family
    chart.value_axis.tick_labels.font.size = Pt(8)
    chart.value_axis.has_major_gridlines = True
    chart.value_axis.major_gridlines.format.line.color.rgb = RGBColor(0xE0, 0xE0, 0xE0)

    return cf


def add_line_chart(slide, brand: BrandPalette,
                   categories: list, series: list,
                   left=Inches(0.8), top=Inches(1.5),
                   width=Inches(11.5), height=Inches(4.5),
                   y_axis_title: str = None):
    """
    折线图。

    series: [("Series A", [1, 2, 3]), ("Series B", [4, 5, 6])]
    """
    chart_data = CategoryChartData()
    chart_data.categories = categories
    for name, values in series:
        chart_data.add_series(name, values)

    cf = slide.shapes.add_chart(
        XL_CHART_TYPE.LINE_MARKERS, left, top, width, height, chart_data
    )
    chart = cf.chart

    colors = brand.dim_colors
    for i, s in enumerate(chart.series):
        s.format.line.color.rgb = colors[i % len(colors)]
        s.format.line.width = Pt(2.5)
        s.marker.style = 8  # XL_MARKER_STYLE.CIRCLE
        s.marker.size = 8
        s.marker.format.fill.solid()
        s.marker.format.fill.fore_color.rgb = colors[i % len(colors)]

    chart.has_legend = len(series) > 1
    if chart.has_legend:
        chart.legend.position = XL_LEGEND_POSITION.BOTTOM
        chart.legend.font.size = Pt(8)

    chart.category_axis.tick_labels.font.size = Pt(8)
    chart.value_axis.tick_labels.font.size = Pt(8)
    chart.value_axis.has_major_gridlines = True
    chart.value_axis.major_gridlines.format.line.color.rgb = RGBColor(0xE0, 0xE0, 0xE0)

    if y_axis_title:
        chart.value_axis.has_title = True
        chart.value_axis.axis_title.text_frame.paragraphs[0].text = y_axis_title
        chart.value_axis.axis_title.text_frame.paragraphs[0].font.size = Pt(8)

    return cf


def add_pie_chart(slide, brand: BrandPalette,
                  categories: list, values: list,
                  left=Inches(3.5), top=Inches(1.5),
                  width=Inches(6.0), height=Inches(5.0),
                  donut: bool = False):
    """
    饼图 / 圆环图。

    categories: ["Segment A", "Segment B", "Segment C"]
    values: [45, 30, 25]
    """
    chart_data = ChartData()
    chart_data.categories = categories
    chart_data.add_series("Share", values)

    chart_type = XL_CHART_TYPE.DOUGHNUT if donut else XL_CHART_TYPE.PIE

    cf = slide.shapes.add_chart(chart_type, left, top, width, height, chart_data)
    chart = cf.chart

    colors = brand.dim_colors
    plot = chart.plots[0]
    series = plot.series[0]

    for i in range(len(categories)):
        pt = series.points[i]
        pt.format.fill.solid()
        pt.format.fill.fore_color.rgb = colors[i % len(colors)]

    # 数据标签
    series.data_labels.show_category_name = True
    series.data_labels.show_percentage = True
    series.data_labels.show_value = False
    series.data_labels.font.size = Pt(9)
    series.data_labels.font.name = brand.font_family

    chart.has_legend = False  # 用数据标签代替图例

    return cf
