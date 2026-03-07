"""
Slide Builders — Reusable McKinsey-style slide templates
========================================================
Each method adds one slide to the presentation following strict design rules.
"""

from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.enum.shapes import MSO_SHAPE
from pptx.enum.chart import XL_CHART_TYPE, XL_LEGEND_POSITION, XL_LABEL_POSITION
from pptx.chart.data import CategoryChartData
from pptx.dml.color import RGBColor

from .design_system import BrandPalette, FONT

# Slide dimensions (16:9 widescreen)
SLIDE_W = Inches(13.333)
SLIDE_H = Inches(7.5)


class DeckBuilder:
    """
    High-level deck builder with McKinsey-style slide templates.

    Usage:
        brand = BrandPalette("#00377B", "#009FDB", "#FFD100")
        deck = DeckBuilder(brand)
        deck.add_cover("Title", "Subtitle")
        deck.add_kpi_slide(...)
        deck.save("output.pptx")
    """

    def __init__(self, brand: BrandPalette):
        self.brand = brand
        self.prs = Presentation()
        self.prs.slide_width = SLIDE_W
        self.prs.slide_height = SLIDE_H

    def _blank(self):
        return self.prs.slides.add_slide(self.prs.slide_layouts[6])

    def _top_bar(self, slide):
        s = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, 0, 0, SLIDE_W, Inches(0.08))
        s.fill.solid(); s.fill.fore_color.rgb = self.brand.primary; s.line.fill.background()

    def _accent_line(self, slide):
        s = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, 0, Inches(0.08), SLIDE_W, Inches(0.03))
        s.fill.solid(); s.fill.fore_color.rgb = self.brand.accent; s.line.fill.background()

    def _footer(self, slide, text="Confidential"):
        bar = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, 0, Inches(7.1), SLIDE_W, Inches(0.4))
        bar.fill.solid(); bar.fill.fore_color.rgb = self.brand.primary; bar.line.fill.background()
        p = bar.text_frame.paragraphs[0]
        p.text = text; p.font.size = Pt(8); p.font.color.rgb = self.brand.white
        p.font.name = self.brand.font_family; p.alignment = PP_ALIGN.RIGHT

    def _action_title(self, slide, title, subtitle=None):
        """Add Action Title (conclusion sentence) + optional subtitle."""
        self._top_bar(slide); self._accent_line(slide)
        tx = slide.shapes.add_textbox(Inches(0.8), Inches(0.25), Inches(11.5), Inches(0.7))
        tf = tx.text_frame; tf.word_wrap = True
        p = tf.paragraphs[0]
        p.text = title; p.font.size = FONT['title']; p.font.bold = True
        p.font.color.rgb = self.brand.primary; p.font.name = self.brand.font_family
        if subtitle:
            p2 = tf.add_paragraph()
            p2.text = subtitle; p2.font.size = FONT['subtitle']
            p2.font.color.rgb = self.brand.dark_gray; p2.font.name = self.brand.font_family
        self._footer(slide)

    def _add_source(self, slide, text, left=Inches(0.8), top=Inches(6.8)):
        src = slide.shapes.add_textbox(left, top, Inches(10), Inches(0.3))
        p = src.text_frame.paragraphs[0]
        p.text = text; p.font.size = FONT['caption']
        p.font.color.rgb = self.brand.mid_gray; p.font.name = self.brand.font_family

    # ═══════════════════════════════════════════
    # SLIDE TEMPLATES
    # ═══════════════════════════════════════════

    def add_cover(self, title: str, subtitle: str, date: str = ""):
        """Full dark background cover slide."""
        slide = self._blank()
        bg = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, 0, 0, SLIDE_W, SLIDE_H)
        bg.fill.solid(); bg.fill.fore_color.rgb = self.brand.primary; bg.line.fill.background()

        # Accent stripe
        s = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(0.8), Inches(2.8), Inches(1.5), Inches(0.06))
        s.fill.solid(); s.fill.fore_color.rgb = self.brand.accent; s.line.fill.background()

        # Title
        tx = slide.shapes.add_textbox(Inches(0.8), Inches(3.0), Inches(10), Inches(1.5))
        tf = tx.text_frame; tf.word_wrap = True
        p = tf.paragraphs[0]
        p.text = title; p.font.size = FONT['cover_title']; p.font.bold = True
        p.font.color.rgb = self.brand.white; p.font.name = self.brand.font_family

        p2 = tf.add_paragraph()
        p2.text = subtitle; p2.font.size = FONT['cover_sub']
        p2.font.color.rgb = self.brand.secondary; p2.font.name = self.brand.font_family
        p2.space_before = Pt(12)

        if date:
            tx2 = slide.shapes.add_textbox(Inches(0.8), Inches(5.5), Inches(10), Inches(0.5))
            p3 = tx2.text_frame.paragraphs[0]
            p3.text = f"{date}  |  CONFIDENTIAL"; p3.font.size = Pt(11)
            p3.font.color.rgb = self.brand.mid_gray; p3.font.name = self.brand.font_family

        # Right vertical stripe
        vs = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(12.5), 0, Inches(0.06), SLIDE_H)
        vs.fill.solid(); vs.fill.fore_color.rgb = self.brand.secondary; vs.line.fill.background()
        return slide

    def add_section_divider(self, title: str, number: str = None):
        """Full dark section divider."""
        slide = self._blank()
        bg = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, 0, 0, SLIDE_W, SLIDE_H)
        bg.fill.solid(); bg.fill.fore_color.rgb = self.brand.primary; bg.line.fill.background()

        s = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(0.8), Inches(3.2), Inches(2.0), Inches(0.05))
        s.fill.solid(); s.fill.fore_color.rgb = self.brand.accent; s.line.fill.background()

        if number:
            tx = slide.shapes.add_textbox(Inches(0.8), Inches(2.0), Inches(2), Inches(1.0))
            p = tx.text_frame.paragraphs[0]
            p.text = f"0{number}"; p.font.size = FONT['section_num']; p.font.bold = True
            p.font.color.rgb = self.brand.accent; p.font.name = self.brand.font_family

        tx2 = slide.shapes.add_textbox(Inches(0.8), Inches(3.5), Inches(10), Inches(1.5))
        p2 = tx2.text_frame.paragraphs[0]
        p2.text = title; p2.font.size = FONT['section_title']; p2.font.bold = True
        p2.font.color.rgb = self.brand.white; p2.font.name = self.brand.font_family
        return slide

    def add_closing(self, title: str = "Thank You", subtitle: str = ""):
        """Closing slide."""
        slide = self._blank()
        bg = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, 0, 0, SLIDE_W, SLIDE_H)
        bg.fill.solid(); bg.fill.fore_color.rgb = self.brand.primary; bg.line.fill.background()

        s = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(4.5), Inches(2.5), Inches(4.0), Inches(0.05))
        s.fill.solid(); s.fill.fore_color.rgb = self.brand.accent; s.line.fill.background()

        tx = slide.shapes.add_textbox(Inches(2), Inches(2.8), Inches(9), Inches(1.5))
        p = tx.text_frame.paragraphs[0]
        p.text = title; p.font.size = FONT['closing']; p.font.bold = True
        p.font.color.rgb = self.brand.white; p.font.name = self.brand.font_family
        p.alignment = PP_ALIGN.CENTER

        if subtitle:
            tx2 = slide.shapes.add_textbox(Inches(2), Inches(4.0), Inches(9), Inches(1.0))
            p2 = tx2.text_frame.paragraphs[0]
            p2.text = subtitle; p2.font.size = Pt(14)
            p2.font.color.rgb = self.brand.secondary; p2.font.name = self.brand.font_family
            p2.alignment = PP_ALIGN.CENTER
        return slide

    def add_kpi_box(self, slide, value: str, label: str,
                    left, top, width=Inches(2.8), height=Inches(1.4),
                    value_color=None, bg_color=None):
        """Large KPI metric box."""
        bg_color = bg_color or self.brand.primary
        value_color = value_color or self.brand.text_on(bg_color)

        box = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, left, top, width, height)
        box.fill.solid(); box.fill.fore_color.rgb = bg_color; box.line.fill.background()

        tf = box.text_frame; tf.word_wrap = True; tf.vertical_anchor = MSO_ANCHOR.MIDDLE
        p1 = tf.paragraphs[0]
        p1.text = value; p1.font.size = FONT['kpi_value']; p1.font.bold = True
        p1.font.color.rgb = value_color; p1.font.name = self.brand.font_family
        p1.alignment = PP_ALIGN.CENTER

        p2 = tf.add_paragraph()
        p2.text = label; p2.font.size = FONT['kpi_label']
        p2.font.color.rgb = self.brand.label_on(bg_color)
        p2.font.name = self.brand.font_family; p2.alignment = PP_ALIGN.CENTER
        return box

    def add_kpi_slide(self, title: str, kpis: list, subtitle: str = None,
                      source: str = None):
        """
        Slide with Action Title + row of KPI boxes.
        kpis: list of (value, label, sublabel) tuples, max 4.
        """
        slide = self._blank()
        self._action_title(slide, title, subtitle)

        n = len(kpis)
        box_w = min(3.0, 11.5 / n - 0.2)
        colors = [self.brand.secondary, self.brand.shade_50,
                  self.brand.primary, self.brand.shade_90]

        for i, kpi in enumerate(kpis):
            val, lbl = kpi[0], kpi[1]
            sublbl = kpi[2] if len(kpi) > 2 else ""
            full_label = f"{lbl}\n{sublbl}" if sublbl else lbl
            x = Inches(0.8 + i * (box_w + 0.15) * (13.333 - 1.6) / (n * box_w + (n-1) * 0.15) * box_w)
            x = Inches(0.8 + i * ((11.5 / n)))
            self.add_kpi_box(slide, val, full_label,
                             x, Inches(1.5), Inches(box_w), Inches(1.4),
                             bg_color=colors[i % len(colors)])

        if source:
            self._add_source(slide, source)
        return slide

    def add_insight_box(self, slide, text: str, left, top,
                        width=Inches(3.5), height=Inches(1.2)):
        """McKinsey-style insight callout box with label bar."""
        box = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, left, top, width, height)
        box.fill.solid(); box.fill.fore_color.rgb = self.brand.tint_05
        box.line.color.rgb = self.brand.secondary; box.line.width = Pt(1.5)

        label = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, left, top, width, Inches(0.28))
        label.fill.solid(); label.fill.fore_color.rgb = self.brand.primary
        label.line.fill.background()
        lp = label.text_frame.paragraphs[0]
        lp.text = "KEY INSIGHT"; lp.font.size = Pt(8); lp.font.bold = True
        lp.font.color.rgb = self.brand.accent; lp.font.name = self.brand.font_family
        lp.alignment = PP_ALIGN.CENTER; label.text_frame.vertical_anchor = MSO_ANCHOR.MIDDLE

        ctf = slide.shapes.add_textbox(left + Inches(0.15), top + Inches(0.35),
                                       width - Inches(0.3), height - Inches(0.45))
        ctf.text_frame.word_wrap = True
        p = ctf.text_frame.paragraphs[0]
        p.text = text; p.font.size = Pt(10)
        p.font.color.rgb = self.brand.charcoal; p.font.name = self.brand.font_family
        return box

    def add_chart_slide(self, title: str, chart_type: str,
                        categories: list, series: list,
                        subtitle: str = None, y_axis_title: str = None,
                        insight: str = None, source: str = None,
                        number_format: str = '$#,##0.0"M"'):
        """
        Slide with Action Title + chart + optional insight box.

        chart_type: 'bar', 'stacked_bar', 'area'
        series: list of (name, values) tuples
        """
        slide = self._blank()
        self._action_title(slide, title, subtitle)

        chart_w = Inches(7.5) if insight else Inches(11.5)
        chart_h = Inches(4.5)
        chart_left = Inches(0.8)
        chart_top = Inches(1.5)

        chart_data = CategoryChartData()
        chart_data.categories = categories
        for name, values in series:
            chart_data.add_series(name, values)

        type_map = {
            'bar': XL_CHART_TYPE.COLUMN_CLUSTERED,
            'stacked_bar': XL_CHART_TYPE.COLUMN_STACKED,
            'area': XL_CHART_TYPE.AREA,
        }
        xl_type = type_map.get(chart_type, XL_CHART_TYPE.COLUMN_CLUSTERED)

        cf = slide.shapes.add_chart(xl_type, chart_left, chart_top, chart_w, chart_h, chart_data)
        chart = cf.chart

        # Legend
        chart.has_legend = len(series) > 1
        if chart.has_legend:
            chart.legend.position = XL_LEGEND_POSITION.BOTTOM
            chart.legend.include_in_layout = False
            chart.legend.font.size = Pt(8); chart.legend.font.name = self.brand.font_family

        # Axes
        chart.category_axis.tick_labels.font.size = Pt(8)
        chart.category_axis.tick_labels.font.name = self.brand.font_family
        chart.value_axis.tick_labels.font.size = Pt(8)
        chart.value_axis.tick_labels.font.name = self.brand.font_family
        chart.value_axis.has_major_gridlines = True
        chart.value_axis.major_gridlines.format.line.color.rgb = RGBColor(0xE0, 0xE0, 0xE0)

        if y_axis_title:
            chart.value_axis.has_title = True
            chart.value_axis.axis_title.text_frame.paragraphs[0].text = y_axis_title
            chart.value_axis.axis_title.text_frame.paragraphs[0].font.size = Pt(8)

        # Series colors + data labels
        colors = self.brand.dim_colors
        for i, s in enumerate(chart.series):
            s.format.fill.solid()
            s.format.fill.fore_color.rgb = colors[i % len(colors)]
            if chart_type != 'stacked_bar' or len(series) <= 2:
                s.data_labels.font.size = Pt(8)
                s.data_labels.font.name = self.brand.font_family
                s.data_labels.show_value = True
                s.data_labels.number_format = number_format

        if insight:
            self.add_insight_box(slide, insight,
                                 Inches(8.8), Inches(1.8), Inches(4.0), Inches(2.0))

        if source:
            self._add_source(slide, source)
        return slide

    def add_table_slide(self, title: str, headers: list, rows: list,
                        subtitle: str = None, source: str = None):
        """Full-width table slide."""
        slide = self._blank()
        self._action_title(slide, title, subtitle)

        n_rows = len(rows) + 1
        n_cols = len(headers)
        table_w = Inches(11.5)
        col_w = table_w // n_cols
        row_h = Inches(0.35)

        ts = slide.shapes.add_table(n_rows, n_cols,
                                    Inches(0.8), Inches(1.3), table_w, row_h * n_rows)
        table = ts.table

        for i, h in enumerate(headers):
            cell = table.cell(0, i)
            cell.text = h; cell.fill.solid(); cell.fill.fore_color.rgb = self.brand.primary
            for p in cell.text_frame.paragraphs:
                p.font.size = FONT['table_header']; p.font.bold = True
                p.font.color.rgb = self.brand.white; p.font.name = self.brand.font_family
                p.alignment = PP_ALIGN.CENTER
            cell.vertical_anchor = MSO_ANCHOR.MIDDLE

        for r_idx, row_data in enumerate(rows):
            bg = self.brand.white if r_idx % 2 == 0 else self.brand.light_gray
            for c_idx, val in enumerate(row_data):
                cell = table.cell(r_idx + 1, c_idx)
                cell.text = str(val); cell.fill.solid(); cell.fill.fore_color.rgb = bg
                for p in cell.text_frame.paragraphs:
                    p.font.size = FONT['table_body']; p.font.color.rgb = self.brand.charcoal
                    p.font.name = self.brand.font_family
                    p.alignment = PP_ALIGN.CENTER if c_idx > 0 else PP_ALIGN.LEFT
                cell.vertical_anchor = MSO_ANCHOR.MIDDLE

        for i in range(n_cols):
            table.columns[i].width = int(col_w)

        if source:
            self._add_source(slide, source)
        return slide

    def add_recommendation_slide(self, title: str, recommendations: list,
                                  subtitle: str = None):
        """
        3-column recommendation cards.
        recommendations: list of (number, title, description) tuples.
        """
        slide = self._blank()
        self._action_title(slide, title, subtitle)

        colors = [self.brand.positive, self.brand.secondary, self.brand.primary]
        n = len(recommendations)

        for i, (num, rec_title, desc) in enumerate(recommendations):
            x = Inches(0.8 + i * (11.5 / n))
            card_w = Inches(11.5 / n - 0.3)
            y = Inches(1.5)

            # Number circle
            circle = slide.shapes.add_shape(MSO_SHAPE.OVAL, x, y, Inches(0.5), Inches(0.5))
            circle.fill.solid(); circle.fill.fore_color.rgb = colors[i % len(colors)]
            circle.line.fill.background()
            cp = circle.text_frame.paragraphs[0]
            cp.text = str(num); cp.font.size = Pt(18); cp.font.bold = True
            cp.font.color.rgb = self.brand.white; cp.font.name = self.brand.font_family
            cp.alignment = PP_ALIGN.CENTER; circle.text_frame.vertical_anchor = MSO_ANCHOR.MIDDLE

            # Title
            ttl = slide.shapes.add_textbox(x + Inches(0.6), y, Inches(card_w - 0.6), Inches(0.5))
            tp = ttl.text_frame.paragraphs[0]
            tp.text = rec_title; tp.font.size = Pt(14); tp.font.bold = True
            tp.font.color.rgb = colors[i % len(colors)]; tp.font.name = self.brand.font_family

            # Description box
            desc_box = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, x, y + Inches(0.6),
                                               card_w, Inches(1.8))
            desc_box.fill.solid(); desc_box.fill.fore_color.rgb = self.brand.light_gray
            desc_box.line.fill.background(); desc_box.text_frame.word_wrap = True
            dp = desc_box.text_frame.paragraphs[0]
            dp.text = desc; dp.font.size = Pt(10)
            dp.font.color.rgb = self.brand.charcoal; dp.font.name = self.brand.font_family
            desc_box.text_frame.vertical_anchor = MSO_ANCHOR.MIDDLE
        return slide

    def save(self, path: str):
        """Save the presentation."""
        self.prs.save(path)
        return len(self.prs.slides)
