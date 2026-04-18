"""
Orchestrator — 端到端编排：采访 → 大纲 → Plan JSON → PPTX
==========================================================
把状态机、Gate、渲染器串成一条可执行的流水线。

用法：
    from engine.orchestrator import Orchestrator
    orc = Orchestrator()
    orc.run_from_interview(interview_answers)  # 自动走完全流程
    # 或
    orc.run_from_outline(outline_data)         # 从大纲开始
    # 或
    orc.run_from_plan("plan.json")             # 从 Plan JSON 开始
"""

import json
from pathlib import Path
from datetime import datetime
from typing import Optional, Dict, Any, List
from enum import Enum

from .state_machine import RunManager, Stage
from .planning_schema import validate_and_report
from .gates import check_pyramid_principle
from .renderer import PlanRenderer
from .design_system import BrandPalette


class FailurePolicy(Enum):
    RETRY = "retry"
    ROLLBACK = "rollback"
    BLOCKED = "blocked"


class Orchestrator:
    """端到端 PPT 生成编排器"""

    DEFAULT_BRAND = {
        "primary": "#00377B",
        "secondary": "#009FDB",
        "accent": "#FFD100",
    }

    def __init__(self, run_id: Optional[str] = None):
        if run_id:
            self.run = RunManager.resume(run_id)
        else:
            self.run = None

    # ================================================================
    # 工具方法：失败处理和环境检查
    # ================================================================

    def _handle_failure(self, stage_name: str, gate_result, attempt: int = 1, max_retries: int = 2):
        """处理 Gate 失败：重试或阻塞"""
        if attempt <= max_retries:
            print(f"  🔄 {stage_name} Gate 失败 (尝试 {attempt}/{max_retries})，重试...")
            return FailurePolicy.RETRY
        else:
            print(f"  ⛔ {stage_name} 连续 {max_retries} 次失败")
            for err in gate_result.errors:
                print(f"     ❌ {err}")
            return FailurePolicy.BLOCKED

    def _preflight_check(self) -> bool:
        """环境预检查：依赖 + 磁盘空间"""
        issues = []
        try:
            import pptx
        except ImportError:
            issues.append("python-pptx 未安装 (pip install python-pptx)")

        import shutil
        free_mb = shutil.disk_usage("/").free / (1024**2)
        if free_mb < 50:
            issues.append(f"磁盘空间不足: {free_mb:.0f}MB (需要 50MB+)")

        if issues:
            print("❌ 环境检查失败:")
            for i in issues:
                print(f"  ⚠️ {i}")
            return False
        print("✅ 环境检查通过")
        return True

    # ================================================================
    # 入口 1：从采访开始（完整流程）
    # ================================================================
    def run_from_interview(self, interview: dict) -> str:
        """
        从采访答案开始，自动走完全流程。

        Args:
            interview: 麦肯锡 10 问的回答（dict）

        Returns:
            输出 PPTX 路径
        """
        if not self._preflight_check():
            return None

        title = interview.get("core_question", "Untitled Deck")
        self.run = RunManager.create(title)

        # P0: 保存采访
        self.run.save_artifact(Stage.INTERVIEW, interview)
        for attempt in range(1, 4):
            gate = self.run.check_gate(Stage.INTERVIEW)
            if gate:
                break
            policy = self._handle_failure("INTERVIEW", gate, attempt)
            if policy == FailurePolicy.BLOCKED:
                print("  💡 请修复 interview 数据后调用 RunManager.resume() 继续")
                return None

        # P1: 自动确认（编排模式下跳过人工确认）
        self.run.save_artifact(Stage.CONFIRM, {
            "confirmed": True,
            "mode": "orchestrated",
            "timestamp": datetime.now().isoformat(),
        })

        # P3: 生成大纲
        outline = self._generate_outline(interview)
        self.run.save_artifact(Stage.OUTLINE, outline)
        for attempt in range(1, 4):
            gate = self.run.check_gate(Stage.OUTLINE)
            if gate:
                break
            policy = self._handle_failure("OUTLINE", gate, attempt)
            if policy == FailurePolicy.BLOCKED:
                print("  💡 请修复 outline 数据后调用 RunManager.resume() 继续")
                return None

        # P3.5: 锁定风格
        brand = interview.get("brand", self.DEFAULT_BRAND)
        style = {"brand": brand, "font_family": "Arial", "locked_at": datetime.now().isoformat()}
        self.run.save_artifact(Stage.STYLE_LOCK, style)

        # P4: 生成 Plan JSON
        for attempt in range(1, 4):
            plan = self._generate_plan(interview, outline, brand)
            self.run.save_artifact(Stage.GENERATE, plan)
            gate = self.run.check_gate(Stage.GENERATE)
            if gate:
                break
            policy = self._handle_failure("GENERATE", gate, attempt)
            if policy == FailurePolicy.BLOCKED:
                print("  💡 请手动编辑 plan.json 修复问题后调用 RunManager.resume() 继续")
                return None

        # P5: 渲染 + QA
        output_path = self._render_and_qa(plan)

        print(f"\n{'='*60}")
        print(f"✅ 全流程完成!")
        print(self.run.status())
        return output_path

    # ================================================================
    # 入口 2：从大纲开始
    # ================================================================
    def run_from_outline(self, outline: dict, brand: dict = None) -> str:
        """从已有大纲开始"""
        if not self._preflight_check():
            return None

        self.run = RunManager.create(outline.get("top_conclusion", "Deck"))

        # 跳过采访和确认
        self.run.save_artifact(Stage.INTERVIEW, {"skipped": True})
        self.run.save_artifact(Stage.CONFIRM, {"confirmed": True, "mode": "from_outline"})
        self.run.save_artifact(Stage.OUTLINE, outline)

        brand = brand or self.DEFAULT_BRAND
        style = {"brand": brand, "font_family": "Arial"}
        self.run.save_artifact(Stage.STYLE_LOCK, style)

        interview = {"brand": brand, "core_question": outline.get("top_conclusion", "")}
        plan = self._generate_plan(interview, outline, brand)
        self.run.save_artifact(Stage.GENERATE, plan)

        return self._render_and_qa(plan)

    # ================================================================
    # 入口 3：从 Plan JSON 开始
    # ================================================================
    def run_from_plan(self, plan_path: str) -> str:
        """从已有 Plan JSON 直接渲染"""
        if not self._preflight_check():
            return None

        plan = json.loads(Path(plan_path).read_text())
        self.run = RunManager.create("Direct Render")

        # 跳过前置阶段
        for stage in [Stage.INTERVIEW, Stage.CONFIRM, Stage.OUTLINE, Stage.STYLE_LOCK]:
            self.run.save_artifact(stage, {"skipped": True})
        self.run.save_artifact(Stage.GENERATE, plan)

        return self._render_and_qa(plan)

    # ================================================================
    # 内部方法
    # ================================================================

    def _generate_outline(self, interview: dict) -> dict:
        """从采访答案生成 Pyramid Principle 大纲"""
        hypothesis = interview.get("hypothesis", "")
        arguments = interview.get("key_arguments", [])
        actions = interview.get("actions", [])

        outline = {
            "top_conclusion": hypothesis,
            "narrative_flow": "conclusion_first",
            "arguments": [],
        }

        for i, arg in enumerate(arguments):
            # 支持 string 或 dict 格式
            if isinstance(arg, dict):
                outline["arguments"].append(arg)
            else:
                slides = [{"template": "data_story", "focus": arg}]
                outline["arguments"].append({
                    "title": arg,
                    "slides": slides,
                    "evidence": [],
                })

        # Pyramid Principle 校验
        issues = check_pyramid_principle(outline)
        criticals = [i for i in issues if i[0] == "CRITICAL"]
        if criticals:
            print(f"  ⚠️ 大纲有 {len(criticals)} 个 CRITICAL 问题，尝试自动修复...")
            # 基本修复：确保至少有 2 个论点
            if len(outline["arguments"]) < 2:
                outline["arguments"].append({
                    "title": "Additional analysis supports the conclusion",
                    "slides": [{"template": "data_story", "focus": "supporting data"}],
                    "evidence": ["To be confirmed"],
                })

        # Density targets based on interview preference
        density_bias = interview.get("density_bias", "balanced")
        density_map = {"relaxed": "medium", "balanced": "high", "ultra_dense": "high"}
        content_density = density_map.get(density_bias, "high")

        outline["density_bias"] = density_bias
        outline["density_targets"] = {
            "cover": "low",
            "section_divider": "low",
            "closing": "low",
            "executive_summary": "medium",
            "recommendation": "medium",
            "data_story": content_density,
            "comparison": content_density,
            "framework": content_density,
            "table": content_density,
            "appendix": content_density,
        }

        return outline

    def _generate_plan(self, interview: dict, outline: dict, brand: dict) -> dict:
        """从大纲生成 Planning JSON"""
        hypothesis = outline.get("top_conclusion", "")
        arguments = outline.get("arguments", [])
        actions = interview.get("actions", [])
        page_count = interview.get("page_count", "8-12")
        date = datetime.now().strftime("%B %Y")

        slides = []
        slide_num = 1

        # 1. Cover — 标题+副标题合计 ≤35 中文字
        cover_title = interview.get("core_question", hypothesis)
        cover_subtitle = hypothesis
        # 如果副标题太长，截断
        cn_total = sum(1 for c in (cover_title + cover_subtitle) if ord(c) > 127)
        if cn_total > 35:
            # 标题优先保留完整，副标题截断
            max_sub = max(10, 35 - sum(1 for c in cover_title if ord(c) > 127))
            cover_subtitle = cover_subtitle[:max_sub] + "…"
        slides.append({
            "slide_number": slide_num,
            "template": "cover",
            "title": cover_title,
            "subtitle": cover_subtitle,
            "date": date,
        })
        slide_num += 1

        # 2. Executive Summary — 用采访中的 kpi_data 或从论点提取
        kpis = []
        kpi_data = interview.get("kpi_data", [])
        if kpi_data:
            # 用户提供了结构化 KPI 数据 — 强制长度限制
            for kpi in kpi_data[:4]:
                label = kpi.get("label", "")
                detail = kpi.get("detail", "")
                # KPI 盒子 2.7"×1.4"，36pt 值 + 10pt 标签，中文合计 ≤12 字
                cn = sum(1 for c in (label + detail) if ord(c) > 127)
                if cn > 0 and len(label + detail) > 12:
                    label = label[:8]
                    detail = detail[:4] if detail else ""
                kpis.append({
                    "value": kpi.get("value", ""),
                    "label": label,
                    "detail": detail,
                })
        else:
            # 从论点标题自动提取数字作为 KPI（尽力而为）
            import re
            for i, arg in enumerate(arguments[:4]):
                title = arg.get("title", "")
                # 尝试提取标题中的第一个数字+单位作为 KPI 值
                num_match = re.search(r'(\d+[\.,]?\d*\s*[%$€¥M万亿]+|\$?\d+[\.,]?\d*[MBK]?\b)', title)
                if num_match:
                    kpis.append({
                        "value": num_match.group(0),
                        "label": title[:40],
                        "detail": "",
                    })
                else:
                    kpis.append({
                        "value": f"Point {i+1}",
                        "label": title[:40],
                        "detail": "",
                    })
        slides.append({
            "slide_number": slide_num,
            "template": "executive_summary",
            "action_title": hypothesis,
            "subtitle": "Key findings summary",
            "kpis": kpis,
            "density_label": "medium",
        })
        slide_num += 1

        # 3-N. 每个论点的支撑页面
        for arg in arguments:
            arg_slides = arg.get("slides", [])
            for s in arg_slides:
                template = s.get("template", "data_story")
                slide_data = {
                    "slide_number": slide_num,
                    "template": template,
                    "action_title": arg.get("title", ""),
                    "subtitle": s.get("focus", ""),
                    "density_label": "medium",
                }

                # data_story — 优先用用户提供的 chart_data
                if template == "data_story":
                    chart_data = s.get("chart_data")
                    if chart_data:
                        slide_data["chart"] = chart_data
                    else:
                        # 无数据时标记为需要填充（Gate 会拦截占位数据）
                        slide_data["chart"] = {
                            "type": "bar",
                            "categories": ["[需要填充真实数据]"],
                            "series": [{"name": "Data", "values": [0]}],
                            "y_axis_title": "[单位]",
                        }
                    slide_data["insight"] = s.get("insight", f"[需要填充: {s.get('focus', '')} 的关键洞察]")
                    # source 从 evidence 获取
                    evidence = arg.get("evidence", [])
                    slide_data["source"] = ", ".join(evidence) if evidence else "[需要填充数据来源]"

                # comparison — 优先用用户提供的对比数据
                elif template == "comparison":
                    comp_data = s.get("comparison_data", {})
                    slide_data["left"] = comp_data.get("left", {
                        "label": "[选项 A]",
                        "points": ["[需要填充对比点]"],
                    })
                    slide_data["right"] = comp_data.get("right", {
                        "label": "[选项 B]",
                        "points": ["[需要填充对比点]"],
                    })

                slides.append(slide_data)
                slide_num += 1

        # N+1. Recommendation（如果有 actions）
        if actions:
            recs = []
            for i, action in enumerate(actions[:5]):  # 支持最多 5 个行动建议
                recs.append({
                    "number": str(i + 1),
                    "title": action.split(":")[0] if ":" in action else action[:30],
                    "detail": action,
                })
            slides.append({
                "slide_number": slide_num,
                "template": "recommendation",
                "action_title": f"{len(actions)} recommended actions to move forward",
                "recommendations": recs,
                "density_label": "medium",
            })
            slide_num += 1

        # Closing
        slides.append({
            "slide_number": slide_num,
            "template": "closing",
            "title": "Thank You",
            "subtitle": f"{interview.get('audience', 'Team')} — {date}",
        })

        plan = {
            "metadata": {
                "title": interview.get("core_question", ""),
                "date": date,
                "brand": brand,
                "total_slides": len(slides),
                "generated_by": "consulting-ppt-skill v2.0",
            },
            "slides": slides,
        }
        return plan

    def _render_and_qa(self, plan: dict) -> str:
        """渲染 PPTX + 简单 QA"""
        # 写临时 plan.json
        plan_path = self.run.run_dir / "plan.json"
        plan_path.write_text(json.dumps(plan, indent=2, ensure_ascii=False))

        # 渲染
        output_path = str(self.run.run_dir / "deck.pptx")
        try:
            renderer = PlanRenderer(str(plan_path))
            renderer.render(output_path)
        except Exception as e:
            print(f"  ❌ 渲染失败: {e}")
            return None

        # 简单 QA
        from pptx import Presentation
        prs = Presentation(output_path)
        qa_report = {
            "total_slides": len(prs.slides),
            "critical_count": 0,
            "warning_count": 0,
            "checks": [],
        }

        for i, slide in enumerate(prs.slides):
            n_shapes = len(slide.shapes)
            if n_shapes < 2:
                qa_report["warning_count"] += 1
                qa_report["checks"].append(f"Slide {i+1}: only {n_shapes} shapes (possibly empty)")
            if n_shapes > 20:
                qa_report["warning_count"] += 1
                qa_report["checks"].append(f"Slide {i+1}: {n_shapes} shapes (possibly cluttered)")

        self.run.save_artifact(Stage.QA, qa_report)
        gate = self.run.check_gate(Stage.QA)

        # 导出清单
        manifest = {
            "output_file": output_path,
            "total_slides": len(prs.slides),
            "qa_passed": bool(gate),
            "generated_at": datetime.now().isoformat(),
        }
        self.run.save_artifact(Stage.EXPORT, manifest)

        return output_path


# ====================================================================
# CLI
# ====================================================================
if __name__ == "__main__":
    import sys

    if len(sys.argv) < 2:
        print("Usage:")
        print("  python -m engine.orchestrator demo     — 运行 demo（麦肯锡 10 问 → PPTX）")
        print("  python -m engine.orchestrator plan.json — 从 Plan JSON 渲染")
        sys.exit(0)

    if sys.argv[1] == "demo":
        orc = Orchestrator()
        result = orc.run_from_interview({
            "audience": "CTO + Board",
            "audience_expertise": "high",
            "core_question": "Should we switch RAN vendors or stay with the current one?",
            "hypothesis": "Staying is 22-57% cheaper than switching across all 4 deployment scenarios",
            "key_arguments": [
                "Migration and multi-vendor costs dominate the switching premium",
                "Stay scenario delivers 16 months faster 5G rollout",
                "Three actions capture value without switching risk",
            ],
            "data_sources": [
                {"name": "TCO Model Q1 2026", "confidence": "A"},
                {"name": "Vendor deployment commitments", "confidence": "A"},
            ],
            "actions": [
                "STAY & NEGOTIATE: Leverage RFP pressure for 10-15% better pricing",
                "MODERNIZE: Use savings to fund 5G acceleration",
                "STRATEGIC RESERVE: Maintain vendor diversification for future",
            ],
            "brand": {"primary": "#00377B", "secondary": "#009FDB", "accent": "#FFD100"},
            "page_count": "8-12",
            "format": "PPTX",
        })
        print(f"\n📄 PPTX: {result}")
    else:
        orc = Orchestrator()
        result = orc.run_from_plan(sys.argv[1])
        print(f"\n📄 PPTX: {result}")
