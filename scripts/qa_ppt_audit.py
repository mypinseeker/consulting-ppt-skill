#!/usr/bin/env python3
"""
QA Audit for McKinsey-style PPT decks  v2.0
=============================================
Maps 1:1 to SKILL.md CRITICAL RULES:

  Rule 1: Max 3 Brand Hues — extracts ALL colors, checks hue count
  Rule 2: Action Titles    — titles must be conclusion sentences
  Rule 3: Every Number Has Unit — bare numbers without $, M, %, etc.
  Rule 4: Contrast >= 3:1  — WCAG contrast on all text
  Rule 5: Charts 60-70%    — chart area + data label presence
  Rule 6: Info Density      — 50-70 chars/sq inch on content slides

Plus formatting checks:
  - Font consistency (non-brand font)
  - Tiny text (< 7pt)
  - Missing source/footnote
  - Empty slides
"""

import sys, os, re, colorsys
from pptx import Presentation
from pptx.util import Pt, Emu
from pptx.dml.color import RGBColor
from pptx.enum.dml import MSO_THEME_COLOR


# ═══════════════════════════════════════════════════════════════
# Color helpers
# ═══════════════════════════════════════════════════════════════

def rgb_luminance(r, g, b):
    """Relative luminance (WCAG formula)."""
    def linearize(v):
        v = v / 255.0
        return v / 12.92 if v <= 0.03928 else ((v + 0.055) / 1.055) ** 2.4
    return 0.2126 * linearize(r) + 0.7152 * linearize(g) + 0.0722 * linearize(b)


def contrast_ratio(c1, c2):
    """WCAG contrast ratio between two (r,g,b) tuples."""
    l1 = rgb_luminance(*c1)
    l2 = rgb_luminance(*c2)
    lighter = max(l1, l2)
    darker = min(l1, l2)
    return (lighter + 0.05) / (darker + 0.05)


def get_shape_bg_rgb(shape):
    """Try to extract background color from shape fill."""
    try:
        if shape.fill and shape.fill.type is not None:
            fc = shape.fill.fore_color
            if fc and fc.type is not None:
                rgb = fc.rgb
                return (rgb[0], rgb[1], rgb[2])
    except Exception:
        pass
    return None


def get_font_rgb(font):
    """Extract font color as (r,g,b) tuple."""
    try:
        if font.color and font.color.rgb:
            rgb = font.color.rgb
            return (rgb[0], rgb[1], rgb[2])
    except Exception:
        pass
    return None


def rgb_to_hue_bucket(r, g, b):
    """Convert RGB to a hue bucket (0-11, i.e. 30° segments).
    Returns None for neutrals (saturation < 10% or very dark/light)."""
    h, s, v = colorsys.rgb_to_hsv(r / 255.0, g / 255.0, b / 255.0)
    # Neutral detection: low saturation OR very dark OR very light
    if s < 0.10 or v < 0.08 or (v > 0.92 and s < 0.08):
        return None  # neutral — white, black, gray
    return int(h * 12) % 12  # 30° buckets


# ═══════════════════════════════════════════════════════════════
# Rule 1: Max 3 Brand Hues
# ═══════════════════════════════════════════════════════════════

def collect_all_colors(prs):
    """Extract every explicit RGB color from shapes, fonts, and charts.
    Returns dict {(r,g,b): [slide_nums]}."""
    color_map = {}

    def _record(rgb_tuple, slide_num):
        if rgb_tuple:
            color_map.setdefault(rgb_tuple, set()).add(slide_num)

    for slide_idx, slide in enumerate(prs.slides, 1):
        for shape in slide.shapes:
            # Shape fill
            _record(get_shape_bg_rgb(shape), slide_idx)

            # Font colors
            if shape.has_text_frame:
                for para in shape.text_frame.paragraphs:
                    for run in para.runs:
                        _record(get_font_rgb(run.font), slide_idx)

            # Chart series colors
            if shape.has_chart:
                try:
                    for series in shape.chart.series:
                        fmt = series.format
                        if fmt and fmt.fill and fmt.fill.type is not None:
                            fc = fmt.fill.fore_color
                            if fc and fc.type is not None:
                                rgb = fc.rgb
                                _record((rgb[0], rgb[1], rgb[2]), slide_idx)
                except Exception:
                    pass

    return color_map


def audit_hue_count(color_map):
    """Check if all colors fit within 3 hue families + neutrals.
    Returns (hue_buckets_dict, issues)."""
    hue_buckets = {}  # bucket_id → set of (r,g,b)
    for rgb in color_map:
        bucket = rgb_to_hue_bucket(*rgb)
        if bucket is not None:  # skip neutrals
            hue_buckets.setdefault(bucket, set()).add(rgb)

    issues = []
    # Merge adjacent buckets — colors within 60° (2 buckets) are same hue family
    merged_families = []
    used = set()
    for b in sorted(hue_buckets):
        if b in used:
            continue
        family = {b}
        # Check adjacent buckets (within ±1, wrapping at 12)
        for delta in [-1, 1]:
            neighbor = (b + delta) % 12
            if neighbor in hue_buckets:
                family.add(neighbor)
                used.add(neighbor)
        used.add(b)
        merged_families.append(family)

    num_families = len(merged_families)
    if num_families > 3:
        # Find which hues are extra
        bucket_names = {
            0: "Red", 1: "Orange", 2: "Yellow", 3: "Yellow-Green",
            4: "Green", 5: "Green-Cyan", 6: "Cyan", 7: "Blue-Cyan",
            8: "Blue", 9: "Blue-Violet", 10: "Violet", 11: "Red-Violet",
        }
        # Describe each family
        family_labels = []
        for fam in merged_families:
            names = [bucket_names.get(b, f"Hue-{b}") for b in sorted(fam)]
            family_labels.append("/".join(names))
        extra_count = num_families - 3
        issues.append(Issue("WARNING", 0, "Deck-wide",
            f"Rule 1: {num_families} hue families detected (max 3): "
            f"{', '.join(family_labels)}. "
            f"{extra_count} extra hue(s) violate the 3-hue rule"))

    return hue_buckets, issues


# ═══════════════════════════════════════════════════════════════
# Rule 2: Action Titles
# ═══════════════════════════════════════════════════════════════

# Section dividers / cover / structural slides are exempt
EXEMPT_TITLE_KEYWORDS = [
    'confidential', 'project ', 'appendix', 'agenda', 'contents',
    'look 1', 'look 2', 'look 3', 'look 4', 'look 5',
    'three decisions', 'thank you', 'q&a',
]

# Verbs that signal a conclusion sentence
VERB_PATTERN = re.compile(
    r'\b(is|are|was|were|has|have|had|'
    r'creates?|reveals?|leads?|drives?|shows?|demands?|requires?|'
    r'generates?|operates?|dominates?|trails?|compounds?|threatens?|'
    r'represents?|faces?|signals?|challenges?|remains?|grows?|'
    r'declines?|exceeds?|achieves?|delivers?|enables?|reduces?|'
    r'increases?|costs?|saves?|limits?|prevents?|outperforms?|'
    r'underperforms?|accelerates?|constrains?|occupies?)\b',
    re.IGNORECASE)


def is_bad_action_title(title_text):
    """Return (is_bad, reason) for a slide title.

    A good McKinsey Action Title:
    - Is a complete sentence stating a conclusion
    - Contains a verb and ideally a number
    - Reads as "so what" not "about what"

    Bad examples: "Market Analysis", "Cost Comparison", "SWOT Matrix"
    Good examples: "Switching costs 22-57% more than staying"
    """
    t = title_text.strip()
    if not t or len(t) < 5:
        return False, ""

    tl = t.lower()

    # Exempt structural slides
    if any(kw in tl for kw in EXEMPT_TITLE_KEYWORDS):
        return False, ""

    words = t.split()
    has_number = bool(re.search(r'\d', t))
    has_verb = bool(VERB_PATTERN.search(t))

    # --- Definite BAD: very short noun phrase, no verb, no number ---
    if len(words) <= 3 and not has_number and not has_verb:
        return True, "too short, no verb, no number — likely a topic label"

    # --- Definite BAD: ends with a broad noun, no verb anywhere ---
    broad_nouns = re.compile(
        r'\b(overview|analysis|comparison|review|landscape|assessment|'
        r'breakdown|composition|update|profile|snapshot|ranking|matrix|'
        r'summary|timeline|introduction|background|discussion|status)$',
        re.IGNORECASE)
    if broad_nouns.search(t) and not has_verb:
        return True, f"ends with topic noun '{broad_nouns.search(t).group()}' and has no verb"

    # --- Suspicious: has a verb but is a template/formulaic filler ---
    # e.g. "Revenue breakdown reveals segment mix and growth drivers"
    # These technically have a verb but the "insight" is generic, not data-specific
    filler_patterns = [
        r'reveals?\s+(key\s+)?(?:strengths|weaknesses|gaps|drivers|trends|insights|factors)',
        r'breakdown\s+reveals?\s+',
        r'assessment\s+reveals?\s+',
        r'check\s+reveals?\s+',
    ]
    for fp in filler_patterns:
        if re.search(fp, t, re.IGNORECASE) and not has_number:
            return True, "verb + generic filler without data — not a real insight"

    # --- Suspicious: ≤ 5 words, no number, no verb ---
    if len(words) <= 5 and not has_number and not has_verb:
        return True, "short title without data or verb"

    # --- OK: has number + verb = likely a real insight ---
    if has_number and has_verb:
        return False, ""

    # --- Marginal: has verb but no number, or vice versa ---
    # Allow but flag for manual review if title is short
    if len(words) <= 6 and not has_number:
        return False, ""  # borderline — don't flag (verb alone is OK for some slides)

    return False, ""


# ═══════════════════════════════════════════════════════════════
# Rule 3: Every Number Has a Unit
# ═══════════════════════════════════════════════════════════════

# Currency prefixes that make a number "have a unit"
CURRENCY_PREFIX_RE = re.compile(
    r'(?:COP|USD|EUR|GBP|BRL|MXN|ARS|PEN|CLP)\s*[\$]?\s*\d', re.IGNORECASE)

# Number followed by a unit suffix
UNIT_SUFFIX_RE = re.compile(
    r'\d\s*(?:%|M\b|K\b|B\b|bn\b|mn\b|k\b|m\b|'
    r'GB|TB|MB|Mbps|Gbps|Kbps|MHz|GHz|'
    r'pp\b|bps\b|pts?\b|x\b|'
    r'months?|years?|days?|hrs?|hours?|'
    r'sites?|bands?|towers?|users?|subs|subscribers?|connections?|'
    r'cities|regions?|provinces?|departments?|'
    r'slides?|scenarios?|steps?|points?|domains?|'
    r'/mo\b|/user|/sub|/month)', re.IGNORECASE)

# Patterns in surrounding text that make bare numbers OK
CONTEXT_UNIT_RE = re.compile(
    r'COP\s|USD\s|\$|€|£|¥|score\s|rating\s|grade\s|rank\s|#\d|'
    r'nPerf|Speedtest|Ookla|Q[1-4]\s|H[12]\s|FY\s|'
    r'Source:|Note:|Confidential|©|March|2026|2025|2024',
    re.IGNORECASE)

# Numbers that are clearly structural (not data)
STRUCTURAL_NUMBER_RE = re.compile(
    r'^\d{1,2}$|'           # single/double digit labels
    r'^0\d$|'               # section numbers "01"
    r'\d+\s+of\s+\d+|'     # "9 of 12" — count expression
    r'#\d+|'               # ranking "#1"
    r'^\d{4}$|'            # year "2025"
    r'\b(?:19|20)\d{2}\b|' # inline year "2025", "2031"
    r'within\s+\d+\s+\w+'  # "within 90 days"
)

# Technology names / labels where digits are part of the name, not data
TECH_LABEL_RE = re.compile(
    r'\b[2-6]G\b|'          # 2G, 3G, 4G, 5G, 6G
    r'\bLTE\b|'
    r'\bWi-?Fi\s*\d|'       # WiFi 6
    r'\bLook\s*\d[\d\-]*|'   # Look 1, Look 2, Looks 1-4
    r'\bDecision\s*\d|'     # Decision 1, Decision 2
    r'\bStep\s*\d|'         # Step 1
    r'\bPhase\s*\d|'        # Phase 1
    r'\bS[1-9]\b|'          # S1, S2 (scenario labels)
    r'\bH[12]\b|'           # H1, H2 (half-year)
    r'\bLooks?\s*\d[\d\-]*|' # Look 1, Looks 1-4
    r'\bP[0-3]\b',          # P0, P1, P2 priority labels
    re.IGNORECASE)

# Count expressions: "N things" where the noun makes the number self-evident
COUNT_NOUN_RE = re.compile(
    r'\b\d+[- ](?:player|operator|factor|metric|dimension|strength|weakness|'
    r'opportunit|threat|domain|pillar|task|scenario|quarter|slide|'
    r'market|country|region|band|item|point|member|day|hour|minute)\w*\b|'
    r'\b\d+\s+(?:SWOT|strengths?|weaknesses?|factors?|metrics?|domains?|'
    r'tasks?|pillars?|priorities?|opportunities)\b|'
    r'(?:across|over|among|between)\s+\d+\s+\w+',   # "across 4 domains"
    re.IGNORECASE)


def has_bare_number(text):
    """Check if text contains numbers without units.
    Significantly improved to reduce false positives."""
    t = text.strip()
    if not t:
        return False

    # Skip text that has context making units implicit
    if CONTEXT_UNIT_RE.search(t):
        return False

    # Skip if text already has currency + number pattern
    if CURRENCY_PREFIX_RE.search(t):
        return False

    # Strip out structural numbers FIRST (years, "X of Y") — before counts break the pattern
    t_clean = STRUCTURAL_NUMBER_RE.sub('__STRUCT__', t)

    # Strip out tech labels (5G, LTE, WiFi 6, etc.) — these are names, not data
    t_clean = TECH_LABEL_RE.sub('__TECH__', t_clean)

    # Strip out count expressions ("4-player", "6 metrics", "12 factors")
    t_clean = COUNT_NOUN_RE.sub('__COUNT__', t_clean)

    # Find remaining numbers after stripping non-data numbers
    remaining_numbers = re.finditer(r'(\d[\d,]*\.?\d*)', t_clean)
    for num_match in remaining_numbers:
        num_str = num_match.group()
        try:
            num_val = float(num_str.replace(',', ''))
        except ValueError:
            continue

        # Skip small numbers (likely counts, indices, scores)
        if num_val < 2.0:
            continue

        # Check 30 chars around the number for unit context
        start = max(0, num_match.start() - 15)
        end = min(len(t_clean), num_match.end() + 15)
        context = t_clean[start:end]

        if UNIT_SUFFIX_RE.search(context):
            continue
        if CURRENCY_PREFIX_RE.search(context):
            continue

        # This number genuinely has no unit
        return True

    return False


# ═══════════════════════════════════════════════════════════════
# Rule 5: Charts Occupy 60-70% of Slide
# ═══════════════════════════════════════════════════════════════

SLIDE_WIDTH_INCHES = 13.333   # default widescreen
SLIDE_HEIGHT_INCHES = 7.5
SLIDE_AREA = SLIDE_WIDTH_INCHES * SLIDE_HEIGHT_INCHES  # ~100 sq in


def chart_area_pct(shape):
    """Return the percentage of slide area occupied by a chart shape."""
    try:
        w = shape.width / 914400.0   # EMU to inches
        h = shape.height / 914400.0
        return (w * h) / SLIDE_AREA * 100
    except Exception:
        return 0


def chart_min_height_inches(shape):
    """Return chart height in inches."""
    try:
        return shape.height / 914400.0
    except Exception:
        return 0


def check_chart_data_labels(chart):
    """Check if chart series have data labels enabled with proper formats.
    Returns list of issues."""
    issues = []
    for s_idx, series in enumerate(chart.series):
        try:
            dl = series.data_labels
            # Check if data labels are shown
            if not dl.show_value and not dl.show_category_name and not dl.show_percentage:
                issues.append(f"Series {s_idx} has no data labels — audience can't read values")
                continue
            # Check format has units
            if dl.show_value:
                fmt = dl.number_format or ''
                if fmt and '$' not in fmt and '%' not in fmt and 'M' not in fmt and '"' not in fmt:
                    issues.append(f"Series {s_idx} data labels format '{fmt}' may lack units")
        except Exception:
            pass
    return issues


# ═══════════════════════════════════════════════════════════════
# Rule 6: Info Density 50-70 chars/sq inch
# ═══════════════════════════════════════════════════════════════

def slide_char_density(slide):
    """Return (total_chars, content_area_sq_in, density_chars_per_sq_in)."""
    total_chars = 0
    for shape in slide.shapes:
        if shape.has_text_frame:
            total_chars += len(shape.text_frame.text.strip())
        # Chart data also contributes to density (implicit content)
        if shape.has_chart:
            total_chars += 100  # charts carry ~100 chars worth of info
    # Content area ≈ 80% of slide (margins excluded)
    content_area = SLIDE_AREA * 0.80
    density = total_chars / content_area if content_area > 0 else 0
    return total_chars, content_area, density


def is_section_divider(slide):
    """Detect section divider slides (full dark background, minimal text)."""
    for shape in slide.shapes:
        bg = get_shape_bg_rgb(shape)
        if bg:
            lum = rgb_luminance(*bg)
            # Dark full-bleed shape covering most of slide
            if lum < 0.15:
                try:
                    w_in = shape.width / 914400.0
                    if w_in > 10:  # wider than 10 inches = full bleed
                        return True
                except Exception:
                    pass
    return False


# ═══════════════════════════════════════════════════════════════
# Issue class + main audit
# ═══════════════════════════════════════════════════════════════

class Issue:
    def __init__(self, severity, slide_num, shape_name, description, text_preview=""):
        self.severity = severity  # CRITICAL / WARNING / INFO
        self.slide_num = slide_num
        self.shape_name = shape_name
        self.description = description
        self.text_preview = text_preview[:80] if text_preview else ""

    def __str__(self):
        preview = f' "{self.text_preview}"' if self.text_preview else ''
        loc = f"Slide {self.slide_num}" if self.slide_num > 0 else "Deck"
        return f"  [{self.severity}] {loc} | {self.shape_name}: {self.description}{preview}"


def audit_pptx(filepath):
    """Run full QA audit on a PPTX file against all 6 SKILL.md rules."""
    prs = Presentation(filepath)
    issues = []
    slide_count = len(prs.slides)

    # ════════════════════════════════════════════
    # RULE 1: Max 3 Brand Hues (deck-wide check)
    # ════════════════════════════════════════════
    color_map = collect_all_colors(prs)
    _, hue_issues = audit_hue_count(color_map)
    issues.extend(hue_issues)

    # ════════════════════════════════════════════
    # Per-slide checks: Rules 2, 3, 4, 5, 6
    # ════════════════════════════════════════════
    for slide_idx, slide in enumerate(prs.slides, 1):
        shape_count = 0
        has_source = False
        divider = is_section_divider(slide)

        for shape in slide.shapes:
            shape_name = shape.name or f"Shape_{shape.shape_id}"
            shape_count += 1

            bg_rgb = get_shape_bg_rgb(shape)

            # ── Text frame checks ──
            if shape.has_text_frame:
                tf = shape.text_frame
                full_text = tf.text.strip()

                if any(kw in full_text.lower() for kw in ['source:', 'note:', 'confidential']):
                    has_source = True

                for para in tf.paragraphs:
                    para_text = para.text.strip()
                    if not para_text:
                        continue

                    font = para.font
                    font_rgb = get_font_rgb(font)
                    font_size = font.size
                    font_name = font.name

                    # RULE 4: Contrast
                    if bg_rgb and font_rgb:
                        ratio = contrast_ratio(font_rgb, bg_rgb)
                        if ratio < 2.0:
                            issues.append(Issue("CRITICAL", slide_idx, shape_name,
                                f"Rule 4: Contrast {ratio:.1f}:1 — text invisible. "
                                f"BG=#{bg_rgb[0]:02X}{bg_rgb[1]:02X}{bg_rgb[2]:02X} "
                                f"Font=#{font_rgb[0]:02X}{font_rgb[1]:02X}{font_rgb[2]:02X}",
                                para_text))
                        elif ratio < 3.0:
                            issues.append(Issue("WARNING", slide_idx, shape_name,
                                f"Rule 4: Contrast {ratio:.1f}:1 < 3.0 minimum. "
                                f"BG=#{bg_rgb[0]:02X}{bg_rgb[1]:02X}{bg_rgb[2]:02X} "
                                f"Font=#{font_rgb[0]:02X}{font_rgb[1]:02X}{font_rgb[2]:02X}",
                                para_text))

                    # RULE 3: Bare numbers
                    if has_bare_number(para_text):
                        issues.append(Issue("WARNING", slide_idx, shape_name,
                            "Rule 3: Number without unit ($M, %, GB, etc.)",
                            para_text))

                    # Font consistency
                    if font_name and font_name not in ('Arial', 'Calibri'):
                        issues.append(Issue("INFO", slide_idx, shape_name,
                            f"Non-standard font: {font_name}",
                            para_text))

                    # Tiny text
                    if font_size and font_size < Pt(7):
                        issues.append(Issue("WARNING", slide_idx, shape_name,
                            f"Text too small: {font_size.pt:.0f}pt (min 7pt)",
                            para_text))

            # ── RULE 5: Chart checks ──
            if shape.has_chart:
                chart = shape.chart
                pct = chart_area_pct(shape)
                h_in = chart_min_height_inches(shape)

                # Chart area
                if pct < 30:
                    issues.append(Issue("WARNING", slide_idx, shape_name,
                        f"Rule 5: Chart occupies only {pct:.0f}% of slide (target 60-70%)"))
                elif pct < 50:
                    issues.append(Issue("INFO", slide_idx, shape_name,
                        f"Rule 5: Chart at {pct:.0f}% — below 60% target"))

                # Chart height
                if h_in < 3.0 and h_in > 0:
                    issues.append(Issue("WARNING", slide_idx, shape_name,
                        f"Rule 5: Chart height {h_in:.1f}\" < 3.0\" minimum"))

                # Data labels
                dl_issues = check_chart_data_labels(chart)
                for dl_msg in dl_issues:
                    issues.append(Issue("WARNING", slide_idx, shape_name,
                        f"Rule 5: {dl_msg}"))

        # ── Empty slide ──
        if shape_count < 2:
            issues.append(Issue("WARNING", slide_idx, "-",
                f"Slide has only {shape_count} shape(s) — possibly empty"))

        # ── Missing source ──
        if slide_idx not in (1, slide_count) and not has_source and not divider:
            issues.append(Issue("INFO", slide_idx, "-",
                "No source/footnote on content slide"))

        # ── RULE 2: Action Titles ──
        # Skip cover (slide 1), last slide, section dividers
        if slide_idx not in (1, slide_count) and not divider:
            title_shape = slide.shapes.title
            if title_shape and title_shape.has_text_frame:
                title_text = title_shape.text_frame.text.strip()
                if title_text:
                    is_bad, reason = is_bad_action_title(title_text)
                    if is_bad:
                        issues.append(Issue("WARNING", slide_idx,
                            title_shape.name or "Title",
                            f"Rule 2: Title is topic, not conclusion — {reason}",
                            title_text))

        # ── RULE 6: Info Density ──
        # Skip cover, closing, dividers
        if slide_idx not in (1, slide_count) and not divider and shape_count >= 2:
            total_chars, content_area, density = slide_char_density(slide)
            if density < 20 and total_chars > 0:
                issues.append(Issue("WARNING", slide_idx, "-",
                    f"Rule 6: Info density {density:.0f} chars/sq in — "
                    f"far below 50-70 target. Add insight boxes or annotations"))
            elif density < 40:
                issues.append(Issue("INFO", slide_idx, "-",
                    f"Rule 6: Info density {density:.0f} chars/sq in — "
                    f"below 50 target, consider adding content"))
            elif density > 100:
                issues.append(Issue("WARNING", slide_idx, "-",
                    f"Rule 6: Info density {density:.0f} chars/sq in — "
                    f"above 70 target, consider splitting"))

    return issues, slide_count


# ═══════════════════════════════════════════════════════════════
# CLI
# ═══════════════════════════════════════════════════════════════

def main():
    if len(sys.argv) > 1:
        ppt_files_full = [f for f in sys.argv[1:] if f.endswith('.pptx')]
    else:
        ppt_files_full = []
        for d in ['.', 'output', 'output/v20260306']:
            full_d = os.path.join(os.path.dirname(__file__), '..', d)
            if os.path.isdir(full_d):
                for f in sorted(os.listdir(full_d)):
                    if f.endswith('.pptx'):
                        ppt_files_full.append(os.path.join(full_d, f))

    if not ppt_files_full:
        print("Usage: python qa_ppt_audit.py [file1.pptx file2.pptx ...]")
        print("       Or place .pptx files in output/ directory")
        return 0

    total_critical = 0
    total_warning = 0
    total_info = 0

    for filepath in ppt_files_full:
        filename = os.path.basename(filepath)
        if not os.path.exists(filepath):
            print(f"\n!! MISSING: {filepath}")
            continue

        issues, slide_count = audit_pptx(filepath)
        criticals = [i for i in issues if i.severity == 'CRITICAL']
        warnings = [i for i in issues if i.severity == 'WARNING']
        infos = [i for i in issues if i.severity == 'INFO']

        total_critical += len(criticals)
        total_warning += len(warnings)
        total_info += len(infos)

        status = "FAIL" if criticals else ("WARN" if warnings else "PASS")
        print(f"\n{'='*70}")
        print(f"  {filename} ({slide_count} slides) — {status}")
        print(f"  CRITICAL: {len(criticals)} | WARNING: {len(warnings)} | INFO: {len(infos)}")
        print(f"{'='*70}")

        # Group by rule for readability
        rule_groups = {}
        other = []
        for issue in criticals + warnings + infos:
            m = re.match(r'Rule (\d)', issue.description)
            if m:
                rule_groups.setdefault(int(m.group(1)), []).append(issue)
            else:
                other.append(issue)

        for rule_num in sorted(rule_groups):
            rule_names = {1: "Max 3 Hues", 2: "Action Titles", 3: "Numbers Have Units",
                          4: "Contrast ≥ 3:1", 5: "Charts 60-70%", 6: "Info Density"}
            print(f"\n  ── Rule {rule_num}: {rule_names.get(rule_num, '?')} ──")
            for issue in rule_groups[rule_num]:
                print(issue)

        if other:
            print(f"\n  ── Formatting ──")
            for issue in other:
                print(issue)

    print(f"\n{'='*70}")
    print(f"  TOTAL: {total_critical} CRITICAL | {total_warning} WARNING | {total_info} INFO")
    gate = "PASS" if total_critical == 0 else "FAIL"
    print(f"  QA GATE: {gate}")
    print(f"{'='*70}")

    return 1 if total_critical > 0 else 0


if __name__ == '__main__':
    sys.exit(main())
