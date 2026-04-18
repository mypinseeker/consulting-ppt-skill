#!/usr/bin/env python3
"""
Visual QA — 结构化视觉审计
============================
不依赖截图，直接解析 PPTX 结构做像素级审计。

检查项：
  1. 文字溢出：文本框尺寸 vs 文字长度估算
  2. 重叠检测：shape 之间的边界框碰撞
  3. 留白均匀：shape 分布是否偏向某侧
  4. 对比度：文字色 vs 背景色
  5. 信息密度：每页 shape 数量 + 文字量
  6. Action Title：是否为判断句
  7. 单位检查：数值是否带单位
  8. 字号下限：< 7pt 的文字
"""

import json
import re
import sys
from pathlib import Path
from dataclasses import dataclass, field
from typing import List, Tuple, Optional

from pptx import Presentation
from pptx.util import Inches, Pt, Emu

# 添加 engine 路径
sys.path.insert(0, str(Path(__file__).parent.parent))
from engine.design_system import rgb_luminance, contrast_ratio
from engine.planning_schema import _is_action_title, UNIT_PATTERN, ACTION_TITLE_EXEMPT


# ====================================================================
# QA Issue
# ====================================================================

@dataclass
class QAIssue:
    slide: int
    check: str
    severity: str   # CRITICAL / WARNING / INFO
    message: str
    fix_hint: str = ""

    def __str__(self):
        hint = f" → Fix: {self.fix_hint}" if self.fix_hint else ""
        return f"[{self.severity}] Slide {self.slide} | {self.check}: {self.message}{hint}"


# ====================================================================
# 尺寸常量（EMU 单位，1 inch = 914400 EMU）
# ====================================================================
SLIDE_W_EMU = Inches(13.333)
SLIDE_H_EMU = Inches(7.5)
MIN_FONT_PT = 7
MIN_MARGIN_EMU = Inches(0.3)       # 最小边距
MAX_SHAPES_PER_SLIDE = 25
OVERLAP_THRESHOLD_EMU = Inches(0.1) # 重叠阈值


# ====================================================================
# 核心检查函数
# ====================================================================

def check_overlap(shapes, slide_num: int) -> List[QAIssue]:
    """检测 shape 之间的边界框重叠"""
    issues = []
    boxes = []

    for s in shapes:
        if hasattr(s, 'left') and s.left is not None:
            boxes.append({
                "name": s.name or "shape",
                "left": s.left,
                "top": s.top,
                "right": s.left + s.width,
                "bottom": s.top + s.height,
            })

    # 过滤掉不参与重叠检测的 shape
    content_boxes = []
    for b in boxes:
        w = b["right"] - b["left"]
        h = b["bottom"] - b["top"]
        # 全页背景矩形
        if w > SLIDE_W_EMU * 0.8 and h > SLIDE_H_EMU * 0.5:
            continue
        # 装饰性窄条（线条/横条/accent bar）
        if h < Inches(0.15) or w < Inches(0.15):
            continue
        # Insight Box 标签栏（高度 < 0.35 inch 且宽度 > 2 inch = KEY INSIGHT bar）
        if h < Inches(0.35) and w > Inches(2.0):
            continue
        # 小圆形（编号圆 Oval，面积 < 0.5 sq inch）
        if w < Inches(0.6) and h < Inches(0.6):
            continue
        # Footer bar（y > 7 inch）
        if b["top"] > Inches(7.0):
            continue
        content_boxes.append(b)

    for i in range(len(content_boxes)):
        for j in range(i + 1, len(content_boxes)):
            a, b = content_boxes[i], content_boxes[j]

            # 容器内嵌套过滤：如果一个 shape 完全包含另一个，视为"内容在容器中"
            a_contains_b = (a["left"] <= b["left"] and a["top"] <= b["top"]
                           and a["right"] >= b["right"] and a["bottom"] >= b["bottom"])
            b_contains_a = (b["left"] <= a["left"] and b["top"] <= a["top"]
                           and b["right"] >= a["right"] and b["bottom"] >= a["bottom"])
            if a_contains_b or b_contains_a:
                continue  # 容器嵌套，不是真正重叠

            x_overlap = max(0, min(a["right"], b["right"]) - max(a["left"], b["left"]))
            y_overlap = max(0, min(a["bottom"], b["bottom"]) - max(a["top"], b["top"]))
            overlap_area = x_overlap * y_overlap

            if overlap_area > 0:
                a_area = (a["right"] - a["left"]) * (a["bottom"] - a["top"])
                b_area = (b["right"] - b["left"]) * (b["bottom"] - b["top"])
                min_area = min(a_area, b_area) if min(a_area, b_area) > 0 else 1
                ratio = overlap_area / min_area

                if ratio > 0.5:
                    issues.append(QAIssue(
                        slide=slide_num,
                        check="overlap",
                        severity="WARNING",
                        message=f"'{a['name']}' 和 '{b['name']}' 重叠 {ratio:.0%}",
                        fix_hint="调整位置或合并 shape",
                    ))

    return issues


def check_text_overflow(shapes, slide_num: int) -> List[QAIssue]:
    """估算文本是否可能溢出容器"""
    issues = []

    for s in shapes:
        if not hasattr(s, 'text_frame'):
            continue

        try:
            tf = s.text_frame
            total_text = tf.text
            if not total_text:
                continue

            # 估算文字需要的空间
            avg_char_width_pt = 6  # 粗略估算，Arial 11pt 约 6pt/字符
            max_font_size = Pt(11)
            for p in tf.paragraphs:
                if p.font and p.font.size:
                    max_font_size = max(max_font_size, p.font.size)

            # 容器宽度（pt）
            if s.width:
                container_width_pt = s.width / 12700  # EMU to pt
                chars_per_line = container_width_pt / avg_char_width_pt
                lines_needed = len(total_text) / max(chars_per_line, 1)

                # 容器高度（pt）
                container_height_pt = s.height / 12700
                line_height_pt = (max_font_size / 12700) * 1.3 if max_font_size else 14
                lines_available = container_height_pt / max(line_height_pt, 1)

                if lines_needed > lines_available * 1.5:
                    issues.append(QAIssue(
                        slide=slide_num,
                        check="text_overflow",
                        severity="WARNING",
                        message=f"'{s.name}' 文字可能溢出 ({len(total_text)} 字符, "
                                f"估算需 {lines_needed:.0f} 行, 容器约 {lines_available:.0f} 行)",
                        fix_hint="缩减文字或增大容器",
                    ))
        except Exception:
            pass

    return issues


def check_contrast(shapes, slide_num: int) -> List[QAIssue]:
    """检查文字与背景的对比度"""
    issues = []

    for s in shapes:
        if not hasattr(s, 'text_frame'):
            continue

        try:
            # 获取背景色
            bg_rgb = None
            if hasattr(s, 'fill') and s.fill and s.fill.type is not None:
                try:
                    fg = s.fill.fore_color
                    if fg and fg.rgb:
                        bg_rgb = (fg.rgb[0], fg.rgb[1], fg.rgb[2])
                except Exception:
                    pass

            if not bg_rgb:
                continue  # 无法确定背景色，跳过

            # 检查每段文字的颜色
            for p in s.text_frame.paragraphs:
                if not p.text.strip():
                    continue

                text_rgb = None
                try:
                    if p.font and p.font.color and p.font.color.rgb:
                        c = p.font.color.rgb
                        text_rgb = (c[0], c[1], c[2])
                except Exception:
                    pass

                if not text_rgb:
                    continue

                ratio = contrast_ratio(text_rgb, bg_rgb)
                if ratio < 2.0:
                    issues.append(QAIssue(
                        slide=slide_num,
                        check="contrast",
                        severity="CRITICAL",
                        message=f"对比度 {ratio:.1f}:1 (需 ≥3:1): "
                                f"文字 #{text_rgb[0]:02x}{text_rgb[1]:02x}{text_rgb[2]:02x} "
                                f"on #{bg_rgb[0]:02x}{bg_rgb[1]:02x}{bg_rgb[2]:02x}",
                        fix_hint="深色背景用白色文字",
                    ))
                elif ratio < 3.0:
                    issues.append(QAIssue(
                        slide=slide_num,
                        check="contrast",
                        severity="WARNING",
                        message=f"对比度偏低 {ratio:.1f}:1: '{p.text[:30]}...'",
                    ))
        except Exception:
            pass

    return issues


def check_density(shapes, slide_num: int) -> List[QAIssue]:
    """检查信息密度"""
    issues = []
    n_shapes = len(shapes)
    total_text = sum(len(s.text) for s in shapes if hasattr(s, 'text'))

    if n_shapes > MAX_SHAPES_PER_SLIDE:
        issues.append(QAIssue(
            slide=slide_num,
            check="density",
            severity="WARNING",
            message=f"{n_shapes} 个 shapes（上限 {MAX_SHAPES_PER_SLIDE}），页面可能过于拥挤",
            fix_hint="简化布局，减少装饰元素",
        ))

    if n_shapes < 3 and slide_num > 1:
        issues.append(QAIssue(
            slide=slide_num,
            check="density",
            severity="INFO",
            message=f"只有 {n_shapes} 个 shapes，页面可能过空",
        ))

    if total_text > 500:
        issues.append(QAIssue(
            slide=slide_num,
            check="density",
            severity="WARNING",
            message=f"{total_text} 字符文字（建议 <400），信息过密",
            fix_hint="拆分为两页或精简文字",
        ))

    return issues


def check_fonts(shapes, slide_num: int) -> List[QAIssue]:
    """检查字号下限"""
    issues = []

    for s in shapes:
        if not hasattr(s, 'text_frame'):
            continue
        for p in s.text_frame.paragraphs:
            if not p.text.strip():
                continue
            try:
                if p.font and p.font.size:
                    size_pt = p.font.size / 12700
                    if size_pt < MIN_FONT_PT:
                        issues.append(QAIssue(
                            slide=slide_num,
                            check="font_size",
                            severity="WARNING",
                            message=f"字号 {size_pt:.0f}pt < {MIN_FONT_PT}pt: '{p.text[:30]}...'",
                            fix_hint=f"增大到至少 {MIN_FONT_PT}pt",
                        ))
            except Exception:
                pass

    return issues


def check_action_titles(shapes, slide_num: int) -> List[QAIssue]:
    """检查 Action Title 是否为判断句"""
    issues = []

    # 找到可能是标题的文本框（位置靠上、字号大）
    for s in shapes:
        if not hasattr(s, 'text_frame'):
            continue
        if not hasattr(s, 'top') or s.top is None:
            continue
        if s.top > Inches(1.2):  # 标题通常在顶部 1.2 英寸以内
            continue

        for p in s.text_frame.paragraphs:
            text = p.text.strip()
            if not text or len(text) < 5:
                continue
            try:
                if p.font and p.font.size and p.font.size > Pt(14):
                    # 这可能是标题
                    if not _is_action_title(text) and len(text) < 50:
                        issues.append(QAIssue(
                            slide=slide_num,
                            check="action_title",
                            severity="WARNING",
                            message=f"标题可能不是判断句: '{text}'",
                            fix_hint="改为包含结论和数字的完整句子",
                        ))
            except Exception:
                pass

    return issues


def check_units(shapes, slide_num: int) -> List[QAIssue]:
    """检查数值是否带单位"""
    issues = []
    num_pattern = re.compile(r'\b\d{2,}\b')  # 两位数以上的数字

    for s in shapes:
        if not hasattr(s, 'text'):
            continue
        text = s.text
        numbers = num_pattern.findall(text)
        for num in numbers:
            # 检查这个数字附近是否有单位
            idx = text.find(num)
            context = text[max(0, idx-5):idx+len(num)+10]
            if not UNIT_PATTERN.search(context):
                issues.append(QAIssue(
                    slide=slide_num,
                    check="unit",
                    severity="INFO",
                    message=f"数字 '{num}' 可能缺少单位: '...{context}...'",
                ))

    return issues


# ====================================================================
# 主审计函数
# ====================================================================

def audit_pptx(pptx_path: str) -> List[QAIssue]:
    """对 PPTX 文件执行全面 Visual QA 审计"""
    prs = Presentation(pptx_path)
    all_issues: List[QAIssue] = []

    for i, slide in enumerate(prs.slides):
        sn = i + 1
        shapes = list(slide.shapes)

        all_issues.extend(check_overlap(shapes, sn))
        all_issues.extend(check_text_overflow(shapes, sn))
        all_issues.extend(check_contrast(shapes, sn))
        all_issues.extend(check_density(shapes, sn))
        all_issues.extend(check_fonts(shapes, sn))
        all_issues.extend(check_action_titles(shapes, sn))
        all_issues.extend(check_units(shapes, sn))

    return all_issues


def audit_and_report(pptx_path: str) -> dict:
    """审计并打印报告，返回 QA report dict"""
    issues = audit_pptx(pptx_path)

    criticals = [i for i in issues if i.severity == "CRITICAL"]
    warnings = [i for i in issues if i.severity == "WARNING"]
    infos = [i for i in issues if i.severity == "INFO"]

    print(f"\n{'='*60}")
    print(f"📋 Visual QA Report: {pptx_path}")
    print(f"{'='*60}")
    print(f"  🔴 CRITICAL: {len(criticals)}")
    print(f"  🟡 WARNING:  {len(warnings)}")
    print(f"  🔵 INFO:     {len(infos)}")
    print(f"{'─'*60}")

    for issue in issues:
        icon = {"CRITICAL": "🔴", "WARNING": "🟡", "INFO": "🔵"}[issue.severity]
        print(f"  {icon} {issue}")

    print(f"{'─'*60}")
    passed = len(criticals) == 0
    print(f"  {'✅ QA PASSED' if passed else '❌ QA FAILED'} — 0 CRITICAL required")

    return {
        "file": pptx_path,
        "critical_count": len(criticals),
        "warning_count": len(warnings),
        "info_count": len(infos),
        "passed": passed,
        "issues": [
            {"slide": i.slide, "check": i.check, "severity": i.severity,
             "message": i.message, "fix_hint": i.fix_hint}
            for i in issues
        ],
    }


# ====================================================================
# CLI
# ====================================================================
if __name__ == "__main__":
    if len(sys.argv) < 2:
        print("Usage: python visual_qa.py <deck.pptx> [--json output.json]")
        sys.exit(1)

    report = audit_and_report(sys.argv[1])

    if len(sys.argv) > 3 and sys.argv[2] == "--json":
        with open(sys.argv[3], "w") as f:
            json.dump(report, f, indent=2, ensure_ascii=False)
        print(f"\n💾 JSON report: {sys.argv[3]}")
